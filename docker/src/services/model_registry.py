"""Model Registry — known model capabilities (context window, vision, limits).

Each provider declares its models with:
  - context_window: max input tokens
  - max_output: max output tokens
  - vision: whether it supports image input
  - file_size_limit_mb: max upload size per file
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Optional


@dataclass
class ModelSpec:
    """Capabilities of a specific model."""
    model_id: str
    provider: str           # openai / anthropic / deepseek / google / ollama / custom
    context_window: int     # max input tokens
    max_output: int         # max output tokens
    vision: bool = False    # supports image input
    max_images: int = 0     # max images per request (0 = N/A)
    file_size_limit_mb: int = 20  # max single file size

    @property
    def safe_chunk_size(self) -> int:
        """40% of context for each chunk."""
        return int(self.context_window * 0.4)

    @property
    def chunk_threshold(self) -> int:
        """If content exceeds 50% of context, trigger chunking."""
        return int(self.context_window * 0.5)


# ── Registry ────────────────────────────────────────────────

KNOWN_MODELS: dict[str, ModelSpec] = {}

def _register(*specs: ModelSpec) -> None:
    for s in specs:
        KNOWN_MODELS[s.model_id] = s


# OpenAI
_register(
    ModelSpec("gpt-4o", "openai", 128000, 16384, vision=True, max_images=20),
    ModelSpec("gpt-4o-mini", "openai", 128000, 16384, vision=True, max_images=20),
    ModelSpec("gpt-4-turbo", "openai", 128000, 4096, vision=True, max_images=10),
    ModelSpec("gpt-4", "openai", 8192, 4096),
    ModelSpec("gpt-3.5-turbo", "openai", 16385, 4096),
    ModelSpec("o3-mini", "openai", 200000, 100000),
    ModelSpec("o1", "openai", 200000, 100000),
)

# Anthropic
_register(
    ModelSpec("claude-sonnet-4-6", "anthropic", 200000, 8192, vision=True, max_images=20),
    ModelSpec("claude-opus-4-7", "anthropic", 200000, 32768, vision=True, max_images=20),
    ModelSpec("claude-haiku-4-5", "anthropic", 200000, 8192, vision=True, max_images=20),
    ModelSpec("claude-3.5-sonnet", "anthropic", 200000, 8192, vision=True, max_images=20),
    ModelSpec("claude-3-opus", "anthropic", 200000, 4096, vision=True, max_images=20),
)

# Kimi / Moonshot
_register(
    ModelSpec("kimi-k2.6", "moonshot", 128000, 8192, vision=True, max_images=10),
    ModelSpec("kimi-k2", "moonshot", 128000, 8192, vision=True, max_images=10),
    ModelSpec("moonshot-v1", "moonshot", 128000, 8192, vision=False),
)

# DeepSeek
_register(
    ModelSpec("deepseek-v4-pro", "deepseek", 1000000, 393216),
    ModelSpec("deepseek-v3", "deepseek", 128000, 8192),
    ModelSpec("deepseek-r1", "deepseek", 128000, 8192),
    ModelSpec("deepseek-chat", "deepseek", 128000, 8192),
)

# Google
_register(
    ModelSpec("gemini-2.5-pro", "google", 1048576, 65536, vision=True, max_images=20),
    ModelSpec("gemini-2.5-flash", "google", 1048576, 65536, vision=True, max_images=20),
    ModelSpec("gemini-2.0-flash", "google", 1048576, 8192, vision=True, max_images=20),
    ModelSpec("gemini-1.5-pro", "google", 2097152, 8192, vision=True, max_images=20),
)

# Alibaba / Qwen
_register(
    ModelSpec("qwen-vl-max", "qwen", 32768, 4096, vision=True, max_images=10),
    ModelSpec("qwen-vl-plus", "qwen", 32768, 4096, vision=True, max_images=10),
    ModelSpec("qwen2.5-vl", "qwen", 32768, 4096, vision=True, max_images=10),
    ModelSpec("qwen3-vl", "qwen", 131072, 8192, vision=True, max_images=10),
    ModelSpec("qwen3.6-plus", "qwen", 131072, 32768, vision=True, max_images=20),
    ModelSpec("qwen3.6", "qwen", 131072, 8192, vision=False),
    ModelSpec("qwen3", "qwen", 32768, 8192, vision=False),
)

# Zhipu / GLM
_register(
    ModelSpec("glm-4v", "zhipu", 128000, 4096, vision=True, max_images=10),
    ModelSpec("glm-4v-plus", "zhipu", 128000, 4096, vision=True, max_images=10),
    ModelSpec("glm-4.6v", "zhipu", 128000, 4096, vision=True, max_images=10),
    ModelSpec("glm-4.5v", "zhipu", 128000, 4096, vision=True, max_images=10),
)

# MiniMax
_register(
    ModelSpec("minimax-m3", "minimax", 128000, 8192, vision=True, max_images=10),
)

# Ollama (conservative defaults, user should check their specific model)
_register(
    ModelSpec("llama3.2-vision", "ollama", 128000, 4096, vision=True, max_images=5),
    ModelSpec("qwen2.5", "ollama", 32768, 4096),
    ModelSpec("llama3.1", "ollama", 128000, 4096),
)


def lookup(model_id: str) -> ModelSpec:
    """Find model spec with fuzzy matching. Exact → prefix → keyword → default."""
    import re
    mid = model_id.lower().strip()

    # 1. Exact match
    if model_id in KNOWN_MODELS:
        return KNOWN_MODELS[model_id]
    if mid in KNOWN_MODELS:
        return KNOWN_MODELS[mid]

    # 2. Strip date suffixes: qwen3.6-plus-2026-03-26 → qwen3.6-plus
    cleaned = re.sub(r'[-_]\d{4}[-_]\d{2}[-_]\d{2}$', '', mid)
    cleaned = re.sub(r'[-_]\d{8}$', '', cleaned)
    cleaned = re.sub(r'[-_]\d{6}$', '', cleaned)
    cleaned = re.sub(r'[-_]20\d{2}$', '', cleaned)
    if cleaned != mid:
        if cleaned in KNOWN_MODELS:
            return KNOWN_MODELS[cleaned]
        for k in KNOWN_MODELS:
            if k.lower() == cleaned:
                return KNOWN_MODELS[k]

    # 3. Prefix matching: find longest matching prefix
    best_match = None
    best_len = 0
    for key, spec in KNOWN_MODELS.items():
        kl = key.lower()
        # Match on first N chars of the key
        min_len = min(len(mid), len(kl))
        # Check if model_id starts with key prefix (e.g., qwen-vl matches qwen-vl-max)
        if mid.startswith(kl) and len(kl) > best_len:
            best_match = spec
            best_len = len(kl)
        # Also check reverse: key starts with model_id
        elif kl.startswith(mid) and len(mid) > best_len:
            best_match = spec
            best_len = len(mid)
    if best_match:
        # Return a copy with the actual model_id
        return ModelSpec(
            model_id=model_id,
            provider=best_match.provider,
            context_window=best_match.context_window,
            max_output=best_match.max_output,
            vision=best_match.vision,
            max_images=best_match.max_images,
            file_size_limit_mb=best_match.file_size_limit_mb,
        )

    # 4. Keyword-based heuristics for unknown models
    vision_keywords = ["vision", "vl", "4o", "claude-3", "claude-4", "sonnet", "opus",
                       "gemini-2", "kimi-k2", "qwen2.5-vl", "glm-4v", "minimax-m3"]
    is_vis = any(kw in mid for kw in vision_keywords)
    large_ctx = any(kw in mid for kw in ["pro", "opus", "turbo", "flash", "k2.6", "v4"])

    return ModelSpec(
        model_id=model_id,
        provider="unknown",
        context_window=128000 if large_ctx else 32768,
        max_output=32768 if large_ctx else 8192,
        vision=is_vis,
    )


def is_vision_model(model_id: str) -> bool:
    spec = lookup(model_id)
    return spec.vision


def get_context_window(model_id: str) -> int:
    return lookup(model_id).context_window


def get_max_output(model_id: str) -> int:
    return lookup(model_id).max_output


def has_images_in_document(content: str, file_ext: str) -> bool:
    """Check if a parsed document likely contains images.

    For PDF/DOCX/PPTX, we need to actually check the source file.
    This function checks the PARSED text for image markers.
    """
    if file_ext in ("pdf", "docx", "pptx", "ppt"):
        # These formats may contain images — need to check source
        # For now, return True as a conservative estimate
        return True
    return False


def extract_images_from_file(file_path: str) -> list[dict]:
    """Extract embedded images from a document. Returns list of {format, width, height, bytes}."""
    ext = file_path.rsplit(".", 1)[-1].lower() if "." in file_path else ""
    images = []

    if ext == "pdf":
        try:
            import fitz
            doc = fitz.open(file_path)
            for page_num, page in enumerate(doc):
                for img in page.get_images(full=True):
                    xref = img[0]
                    try:
                        info = doc.extract_image(xref)
                        images.append({
                            "format": info["ext"],
                            "width": info["width"],
                            "height": info["height"],
                            "page": page_num + 1,
                            "bytes": info["image"],
                        })
                    except Exception:
                        pass
            doc.close()
        except ImportError:
            pass

    elif ext in ("docx",):
        try:
            from docx import Document
            import io
            doc = Document(file_path)
            for rel in doc.part.rels.values():
                if "image" in (rel.reltype or ""):
                    try:
                        img_data = rel.target_part.blob
                        images.append({
                            "format": rel.target_ext or "png",
                            "width": 0, "height": 0,
                            "page": 0,
                            "bytes": img_data,
                        })
                    except Exception:
                        pass
        except ImportError:
            pass

    elif ext in ("pptx", "ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        try:
                            img = shape.image
                            images.append({
                                "format": img.content_type.split("/")[-1] if "/" in (img.content_type or "") else "png",
                                "width": img.size[0] if hasattr(img, 'size') else 0,
                                "height": img.size[1] if hasattr(img, 'size') else 0,
                                "page": slide_num + 1,
                                "bytes": img.blob,
                            })
                        except Exception:
                            pass
        except ImportError:
            pass

    return images


def build_image_section(images: list[dict], model: ModelSpec) -> list[dict]:
    """Build content blocks for a vision model request. Returns Anthropic-style blocks."""
    blocks = []
    for i, img in enumerate(images):
        if i >= model.max_images:
            break
        # Convert to base64
        import base64
        b64 = base64.b64encode(img["bytes"]).decode("utf-8")
        mime = f"image/{img['format']}" if img["format"] else "image/png"
        blocks.append({
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": mime,
                "data": b64,
            },
        })
    return blocks
