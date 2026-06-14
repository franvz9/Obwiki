"""Multi-format document parser. Converts binary documents to Markdown text.

Supported formats (matching llm_wiki's Rust backend):
  PDF  → pymupdf (MuPDF) with optional image extraction
  DOCX → python-docx (headings, bold/italic, lists, tables)
  PPTX → python-pptx (slide-by-slide with heading/list structure)
  XLSX → openpyxl (sheet-by-sheet with Markdown tables)
  XLS  → xlrd (legacy Excel)
  ODS  → pandas/openpyxl
  Images → Pillow metadata
  Text → direct UTF-8 read
  HTML → html2text
  CSV  → pandas summary
"""

from __future__ import annotations
from ..paths import INBOX_DIR, META_DIR, WIKI_DIR, COMMUNITIES_DIR, LLMWIKI_DIR, OPS_DIR, INDEX_FILE, OVERVIEW_FILE, LOG_FILE, SCHEMA_FILE, CACHE_DIR

import hashlib
import os
from io import StringIO
from pathlib import Path
from typing import Optional

# Known extensions mapping
OFFICE_EXTS = {"doc", "docx", "pptx", "xls", "xlsx", "odt", "ods", "odp", "ppt"}
IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "webp", "bmp", "ico", "tiff", "tif", "avif", "heic", "heif", "svg"}
MEDIA_EXTS = {"mp4", "webm", "mov", "avi", "mkv", "mp3", "wav", "ogg", "flac", "aac", "m4a", "wma"}
LEGACY_EXTS = {"pages", "numbers", "key", "epub"}

# CACHE_DIR imported from paths


def get_ext(path: str) -> str:
    return Path(path).suffix.lstrip(".").lower()


def parse_document(path: str) -> str:
    """Parse any supported document to Markdown text. Auto-detects format by extension."""
    p = Path(path)
    if not p.exists():
        return f"[File not found: {path}]"

    ext = get_ext(path)

    # Check cache first
    cached = _read_cache(p)
    if cached is not None:
        return cached

    result: str
    match ext:
        case "pdf":
            result = _parse_pdf(path)
        case "docx":
            result = _parse_docx(path)
        case "doc":
            result = _parse_doc_legacy(path)
        case "pptx" | "ppt":
            result = _parse_pptx(path)
        case "xlsx" | "xls":
            result = _parse_spreadsheet(path, ext)
        case "ods":
            result = _parse_spreadsheet(path, ext)
        case "odt" | "odp":
            result = _parse_odf(path)
        case e if e in IMAGE_EXTS:
            size = p.stat().st_size
            result = f"[Image: {p.name} ({size / 1024:.1f} KB)]"
        case e if e in MEDIA_EXTS:
            size = p.stat().st_size
            result = f"[Media: {p.name} ({size / 1048576:.1f} MB)]"
        case e if e in LEGACY_EXTS:
            result = f"[Document: {p.name} — text extraction not supported for .{e} format]"
        case "html" | "htm":
            result = _parse_html(path)
        case "csv":
            result = _parse_csv(path)
        case "json":
            result = _parse_json_safe(path)
        case _:
            result = _read_text(path)

    # Cache the result
    _write_cache(p, result)
    return result


def _read_text(path: str) -> str:
    try:
        return Path(path).read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return f"[Binary file: {Path(path).name} — content not displayed]"


def _parse_pdf(path: str) -> str:
    try:
        import fitz  # pymupdf
        doc = fitz.open(path)
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text("text")
            if text.strip():
                pages.append(f"## Page {i + 1}\n\n{text}")
        doc.close()
        return "\n\n".join(pages) if pages else "[PDF: no extractable text]"
    except ImportError:
        pass

    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            pages = []
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    pages.append(f"## Page {i + 1}\n\n{text}")
            return "\n\n".join(pages) if pages else "[PDF: no extractable text]"
    except ImportError:
        pass

    return "[PDF: pdfplumber/pymupdf not installed — cannot extract text]"


def _parse_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        result = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                result.append("")
                continue
            if para.style.name.startswith("Heading"):
                level = para.style.name.split()[-1]
                try:
                    lv = int(level)
                    result.append(f"{'#' * lv} {text}")
                except ValueError:
                    result.append(f"## {text}")
            elif para.style.name == "List Bullet":
                result.append(f"- {text}")
            elif para.style.name == "List Number":
                result.append(f"1. {text}")
            else:
                result.append(text)

        # Tables
        for table in doc.tables:
            result.append("")
            rows = []
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                result.append(rows[0])
                result.append("|" + "---|" * len(table.rows[0].cells))
                for r in rows[1:]:
                    result.append(r)

        return "\n".join(result) if result else "[DOCX: no extractable text]"
    except ImportError:
        return "[DOCX: python-docx not installed]"


def _parse_doc_legacy(path: str) -> str:
    ext = get_ext(path)
    if ext != "doc":
        return _read_text(path)

    try:
        # Try antiword or textract
        import subprocess
        result = subprocess.run(
            ["antiword", path], capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    try:
        import textract
        text = textract.process(path).decode("utf-8")
        return text
    except ImportError:
        pass

    return "[DOC: legacy Word format — install antiword or textract for extraction]"


def _parse_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            lines = [f"## Slide {i + 1}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            level = para.level or 0
                            prefix = "  " * level + "- " if level > 0 else ""
                            lines.append(f"{prefix}{text}")
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        rows.append("| " + " | ".join(cells) + " |")
                    if rows:
                        lines.append("")
                        lines.extend(rows)
            slides.append("\n".join(lines))
        return "\n\n".join(slides) if slides else "[PPTX: no extractable text]"
    except ImportError:
        return "[PPTX: python-pptx not installed]"


def _parse_spreadsheet(path: str, ext: str) -> str:
    import csv

    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        sheets = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            sheets.append(f"## Sheet: {name}\n")
            # Max 200 rows per sheet
            for row in rows[:200]:
                cells = [str(c) if c is not None else "" for c in row]
                sheets.append("| " + " | ".join(cells) + " |")
            sheets.append("")
        wb.close()
        return "\n".join(sheets) if sheets else "[XLSX: empty workbook]"
    except ImportError:
        pass

    if ext == "xls":
        try:
            import xlrd
            wb = xlrd.open_workbook(path)
            sheets = []
            for name in wb.sheet_names():
                ws = wb.sheet_by_name(name)
                sheets.append(f"## Sheet: {name}\n")
                for row_idx in range(min(ws.nrows, 200)):
                    cells = [str(ws.cell_value(row_idx, col)) for col in range(ws.ncols)]
                    sheets.append("| " + " | ".join(cells) + " |")
                sheets.append("")
            return "\n".join(sheets) if sheets else "[XLS: empty workbook]"
        except ImportError:
            pass

    return "[Spreadsheet: install openpyxl + xlrd]"


def _parse_odf(path: str) -> str:
    try:
        import zipfile
        from xml.etree import ElementTree as ET

        with zipfile.ZipFile(path) as z:
            if "content.xml" in z.namelist():
                xml = z.read("content.xml")
                root = ET.fromstring(xml)
                ns = {"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
                      "office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0"}
                texts = []
                for elem in root.iter():
                    tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
                    if tag in ("p", "h"):
                        t = "".join(elem.itertext()).strip()
                        if t:
                            prefix = "## " if tag == "h" else ""
                            texts.append(f"{prefix}{t}")
                return "\n\n".join(texts) if texts else "[ODF: no extractable text]"
    except Exception:
        pass
    return "[ODF: extraction failed]"


def _parse_html(path: str) -> str:
    try:
        import html2text
        h = html2text.HTML2Text()
        h.body_width = 0
        h.ignore_links = False
        h.ignore_images = True
        return h.handle(Path(path).read_text(encoding="utf-8"))
    except ImportError:
        pass

    # fallback: simple tag stripping
    import re
    text = Path(path).read_text(encoding="utf-8")
    text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<[^>]+>', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _parse_csv(path: str) -> str:
    try:
        import pandas as pd
        df = pd.read_csv(path, nrows=100)
        return df.to_markdown(index=False)
    except ImportError:
        pass

    import csv
    output = StringIO()
    with open(path) as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            if i > 200:
                break
            output.write("| " + " | ".join(row) + " |\n")
    return output.getvalue()


def _parse_json_safe(path: str) -> str:
    try:
        import json
        data = json.loads(Path(path).read_text(encoding="utf-8"))
        return f"```json\n{json.dumps(data, ensure_ascii=False, indent=2)[:5000]}\n```"
    except Exception:
        pass
    return _read_text(path)


# ── Cache helpers ─────────────────────────────────────────────────

def _cache_path(original: Path) -> Path:
    parent = original.parent
    cache_dir = parent / CACHE_DIR
    return cache_dir / f"{original.name}.txt"


def _read_cache(original: Path) -> Optional[str]:
    cp = _cache_path(original)
    if not cp.exists():
        return None
    if original.stat().st_mtime > cp.stat().st_mtime:
        return None
    try:
        return cp.read_text(encoding="utf-8")
    except Exception:
        return None


def _write_cache(original: Path, text: str) -> None:
    cp = _cache_path(original)
    cp.parent.mkdir(parents=True, exist_ok=True)
    cp.write_text(text, encoding="utf-8")
